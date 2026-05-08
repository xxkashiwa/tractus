from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path


@dataclass(frozen=True)
class DocConvertResult:
    """文档转换为 Markdown 的结果。"""
    markdown_content: str
    source_file: str
    file_type: str
    images_dir: str | None = None
    image_count: int = 0


def convert_doc_to_markdown(
    file_path: str | Path,
    *,
    output_dir: str | Path | None = None,
    remove_images: bool = False,
) -> DocConvertResult:
    """使用 markitdown 将文档文件转换为 Markdown 格式。

    支持 PDF、DOCX、PPTX、HTML 等多种文档格式。
    可选项：从文档中提取图片并保存到 output_dir。

    参数：
        file_path: 要转换的文档文件路径。
        output_dir: 保存提取图片的目录。如果为 None，则不提取图片。
        remove_images: 如果为 True，从输出 Markdown 中移除图片引用。

    返回值：
        包含转换后的 Markdown 内容和元数据的 DocConvertResult。

    异常：
        ValueError: 如果文件路径为空或文件不存在。
        RuntimeError: 如果 markitdown 未安装或转换失败。
        FileNotFoundError: 如果找不到指定的文件。
    """
    raw_file_path = str(file_path).strip()
    if not raw_file_path:
        raise ValueError("file_path must be non-empty")

    file_path = Path(raw_file_path)
    
    if not file_path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")
    
    try:
        import markitdown
    except ImportError:
        raise RuntimeError(
            "markitdown is not installed. Install it with: pip install markitdown"
        )
    
    try:
        # 创建 MarkItDown 转换器
        md = markitdown.MarkItDown()
        
        # 将文档转换为 markdown
        result = md.convert(str(file_path))
        
        # 获取 markdown 内容
        markdown_content = result.text_content
        
        # 如果提供了 output_dir，提取图片
        images_dir = None
        image_count = 0
        
        if output_dir and not remove_images:
            output_path = Path(output_dir)
            images_dir_path = output_path / "images"
            images_dir_path.mkdir(parents=True, exist_ok=True)
            
            # 使用通用方法提取图片
            image_count, page_images_map = _extract_images(
                file_path, images_dir_path, markdown_content, dpi=300
            )
            if image_count > 0:
                # 在适当位置将图片添加到 markdown
                markdown_content = _add_page_markers_and_images(
                    markdown_content, page_images_map
                )
                images_dir = str(images_dir_path)
        
        # 可选项：移除图片引用
        if remove_images:
            markdown_content = _remove_images(markdown_content)
        
        return DocConvertResult(
            markdown_content=markdown_content,
            source_file=str(file_path),
            file_type=file_path.suffix.lower(),
            images_dir=images_dir,
            image_count=image_count,
        )
    
    except Exception as e:
        raise RuntimeError(
            f"Failed to convert document '{file_path}': {str(e)}"
        ) from e


def _extract_images(doc_path: Path, output_dir: Path, markdown_content: str, dpi: int = 300) -> tuple[int, dict]:
    """通用的图片提取方法，支持多种文档格式。
    
    参数：
        doc_path: 文档文件的路径。
        output_dir: 保存提取图片的目录。
        markdown_content: 当前 markdown 内容。
        dpi: 渲染的每英寸点数，用于提高图片质量。
    
    返回值：
        (image_count, page_images_map) 的元组。
    """
    file_ext = doc_path.suffix.lower()
    
    if file_ext == ".pdf":
        return _extract_pdf_images(doc_path, output_dir, markdown_content, dpi=dpi)
    elif file_ext == ".docx":
        return _extract_docx_images(doc_path, output_dir, dpi=dpi)
    elif file_ext == ".pptx":
        return _extract_pptx_images(doc_path, output_dir, dpi=dpi)
    else:
        # 对于不支持的格式，返回空结果
        return 0, {}



def _extract_pdf_images(pdf_path: Path, output_dir: Path, markdown_content: str, dpi: int = 300) -> tuple[int, dict]:
    """以高分辨率按页面从 PDF 文件中提取图片。
    
    参数：
        pdf_path: PDF 文件的路径。
        output_dir: 保存提取图片的目录。
        markdown_content: 当前 markdown 内容（用于检查图片是否被引用）。
        dpi: 渲染的每英寸点数。更高的值（例如 300）提供更好的质量。
    
    返回值：
        (image_count, page_images_map) 的元组，其中 page_images_map 是一个字典
        将页码映射到图片文件名列表。
    """
    try:
        import pdfplumber
    except ImportError:
        return 0, {}
    
    image_count = 0
    page_images_map = {}  # {page_num: [image_filenames]}
    
    # 从 DPI 计算缩放因子（标准为 72 DPI）
    zoom = dpi / 72.0
    
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                page_images = []
                # 从页面中提取图片，使用缩放以获得更高质量
                for img_idx, img in enumerate(page.images):
                    try:
                        # 获取图片边界框并使用缩放裁剪
                        cropped = page.crop((img["x0"], img["top"], img["x1"], img["bottom"]))
                        # 使用分辨率参数以获得更高的分辨率
                        im = cropped.to_image(resolution=dpi)
                        
                        # 以高质量保存图片
                        img_filename = f"page_{page_num}_image_{img_idx}.png"
                        img_path = output_dir / img_filename
                        im.save(img_path, quality=95)
                        
                        page_images.append(img_filename)
                        image_count += 1
                    except Exception as e:
                        print(f"警告：无法从第 {page_num} 页提取图片：{e}")
                        continue
                
                if page_images:
                    page_images_map[page_num] = page_images
    except Exception as e:
        print(f"警告：无法从 PDF 中提取图片：{e}")
    
    return image_count, page_images_map


def _add_page_markers_and_images(markdown_content: str, page_images_map: dict) -> str:
    """为 markdown 添加页面感知的图片引用。
    
    参数：
        markdown_content: 原始 markdown 内容。
        page_images_map: 将页码映射到图片文件名列表的字典。
    
    返回值：
        在适当页面位置嵌入图片的更新后的 markdown。
    """
    if not page_images_map:
        return markdown_content
    
    # 大致按页面分割 markdown
    lines = markdown_content.split('\n')
    num_pages = max(page_images_map.keys()) if page_images_map else 1
    updated_lines = []
    current_page = 1
    target_lines_per_page = max(1, len(lines) // num_pages) if num_pages > 0 else len(lines)
    lines_in_current_page = 0
    
    for line in lines:
        updated_lines.append(line)
        lines_in_current_page += 1
        
        # 检查是否应该转到下一页
        if lines_in_current_page >= target_lines_per_page and current_page < num_pages:
            # 在转到下一页之前为当前页添加图片
            if current_page in page_images_map:
                updated_lines.append("")  # 空行用于间隔
                for img_filename in page_images_map[current_page]:
                    relative_path = f"images/{img_filename}"
                    updated_lines.append(f"![{img_filename.replace('.png', '')}]({relative_path})")
                updated_lines.append("")
            
            current_page += 1
            lines_in_current_page = 0
    
    # 为最后一页添加图片
    if current_page in page_images_map:
        updated_lines.append("")
        for img_filename in page_images_map[current_page]:
            relative_path = f"images/{img_filename}"
            updated_lines.append(f"![{img_filename.replace('.png', '')}]({relative_path})")
        updated_lines.append("")
    
    return '\n'.join(updated_lines)


def _extract_docx_images(docx_path: Path, output_dir: Path, dpi: int = 300) -> tuple[int, dict]:
    """从 DOCX 文件中提取图片。
    
    参数：
        docx_path: DOCX 文件的路径。
        output_dir: 保存提取图片的目录。
        dpi: 分辨率参数（为了保持 API 一致性）。
    
    返回值：
        (image_count, page_images_map) 的元组。
    """
    try:
        from docx import Document
        from docx.oxml import parse_xml
        from docx.oxml.ns import nsdecls
    except ImportError:
        print("警告：python-docx 未安装。请使用 pip install python-docx")
        return 0, {}
    
    image_count = 0
    page_images_map = {}
    
    try:
        doc = Document(docx_path)
        
        # 遍历文档中的所有关系以找到嵌入的图片
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    # 获取图片数据
                    image_data = rel.target_part.blob
                    
                    # 获取图片扩展名
                    target_ref = rel.target_ref
                    img_ext = target_ref.split('.')[-1] if '.' in target_ref else 'png'
                    
                    # 保存图片
                    img_filename = f"docx_image_{image_count}.{img_ext}"
                    img_path = output_dir / img_filename
                    with open(img_path, 'wb') as f:
                        f.write(image_data)
                    
                    image_count += 1
                    # DOCX 没有明确的页数概念，都映射到第 1 页
                    if 1 not in page_images_map:
                        page_images_map[1] = []
                    page_images_map[1].append(img_filename)
                except Exception as e:
                    print(f"警告：无法提取 DOCX 图片：{e}")
                    continue
    except Exception as e:
        print(f"警告：无法处理 DOCX 文件：{e}")
    
    return image_count, page_images_map


def _extract_pptx_images(pptx_path: Path, output_dir: Path, dpi: int = 300) -> tuple[int, dict]:
    """从 PPTX 文件中提取图片。
    
    参数：
        pptx_path: PPTX 文件的路径。
        output_dir: 保存提取图片的目录。
        dpi: 分辨率参数（为了保持 API 一致性）。
    
    返回值：
        (image_count, page_images_map) 的元组，其中 page_images_map 按幻灯片分组。
    """
    try:
        from pptx import Presentation
    except ImportError:
        print("警告：python-pptx 未安装。请使用 pip install python-pptx")
        return 0, {}
    
    image_count = 0
    page_images_map = {}
    
    try:
        prs = Presentation(pptx_path)
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_images = []
            
            # 遍历幻灯片中的所有形状
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        
                        # 获取图片文件名和扩展名
                        img_filename = f"slide_{slide_num}_image_{len(slide_images)}.{image.ext}"
                        img_path = output_dir / img_filename
                        
                        with open(img_path, 'wb') as f:
                            f.write(image_bytes)
                        
                        slide_images.append(img_filename)
                        image_count += 1
                    except Exception as e:
                        print(f"警告：无法从第 {slide_num} 张幻灯片提取图片：{e}")
                        continue
            
            if slide_images:
                page_images_map[slide_num] = slide_images
    except Exception as e:
        print(f"警告：无法处理 PPTX 文件：{e}")
    
    return image_count, page_images_map



    """从 Markdown 内容中移除图片引用。
    
    移除内联图片 ![alt](url) 和引用式图片。
    
    参数：
        markdown_content: 包含潜在图片引用的 Markdown 内容。
    
    返回值：
        移除了图片引用的 Markdown 内容。
    """
    import re
    
    # 移除内联图片：![alt](url)
    markdown_content = re.sub(r'!\[.*?\]\(.*?\)', '', markdown_content)
    
    # 移除引用式图片：![alt][ref]
    markdown_content = re.sub(r'!\[.*?\]\[.*?\]', '', markdown_content)
    
    # 移除图片引用定义：[ref]: url
    markdown_content = re.sub(r'^\[.*?\]:\s*\S+$', '', markdown_content, flags=re.MULTILINE)
    
    return markdown_content


def convert_docs_batch(
    directory: str | Path,
    *,
    output_dir: str | Path | None = None,
    file_extensions: list[str] | None = None,
    remove_images: bool = False,
) -> list[DocConvertResult]:
    """将目录中的多个文档转换为 Markdown 格式。
    
    参数：
        directory: 包含要转换的文档的目录。
        output_dir: 保存转换后的 markdown 文件和图片的目录。
        file_extensions: 要处理的文件扩展名列表（例如 ['.pdf', '.docx']）。
                        如果为 None，则处理常见的文档格式。
        remove_images: 如果为 True，从输出 Markdown 中移除图片引用。
    
    返回值：
        每个成功转换的文件的 DocConvertResult 列表。
    
    异常：
        ValueError: 如果目录路径为空。
        NotADirectoryError: 如果指定的路径不是目录。
        RuntimeError: 如果 markitdown 未安装。
    """
    raw_directory = str(directory).strip()
    if not raw_directory:
        raise ValueError("directory must be non-empty")

    directory = Path(raw_directory)
    
    if not directory.is_dir():
        raise NotADirectoryError(f"Not a directory: {directory}")
    
    if file_extensions is None:
        # markitdown 支持的常见文档格式
        file_extensions = ['.pdf', '.docx', '.pptx', '.html', '.txt', '.md', '.xlsx', '.csv']
    
    results = []
    
    for file_path in directory.rglob('*'):
        if not file_path.is_file():
            continue
        
        if file_path.suffix.lower() not in file_extensions:
            continue
        
        try:
            result = convert_doc_to_markdown(
                file_path,
                output_dir=output_dir,
                remove_images=remove_images,
            )
            results.append(result)
        except Exception as e:
            # 记录错误但继续处理其他文件
            print(f"警告：无法转换 {file_path}：{e}")
            continue
    
    return results
